import fitz  # PyMuPDF
from pdf2image import convert_from_path
import pytesseract
from docx import Document
from pptx import Presentation

# Update path if needed for Windows
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

def extract_text(file_path, file_type):
    if file_type == "pdf":
        return extract_text_from_pdf(file_path)
    elif file_type == "docx":
        return extract_text_from_word(file_path)
    elif file_type == "pptx":
        return extract_text_from_ppt(file_path)
    else:
        return "Unsupported file type."

def extract_text_from_pdf(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        text = ""
        for page in doc:
            page_text = page.get_text()
            if page_text.strip():
                text += page_text
        if text.strip():
            return text
        else:
            raise ValueError("PDF might be scanned. Using OCR...")
    except Exception:
        return extract_text_with_ocr(pdf_path)

def extract_text_with_ocr(pdf_path):
    images = convert_from_path(pdf_path)
    text = ""
    for img in images:
        text += pytesseract.image_to_string(img)
    return text

def extract_text_from_word(docx_path):
    doc = Document(docx_path)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
